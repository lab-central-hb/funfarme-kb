"""
gera-relatorio.py — Gera relatório mensal de atividades

Lê docs/tarefas.md + docs/incidentes/ e gera:
  - portal/docs/relatorios/YYYY-MM.md  (versão web)
  - reports/YYYY-MM_resumo-mensal.pptx (apresentação)
  - Atualiza portal/docs/relatorios/index.md
  - Atualiza contadores no portal/docs/index.md

Uso: python scripts/gera-relatorio.py [--mes YYYY-MM]
"""

import argparse
import re
import os
from datetime import datetime
from collections import defaultdict, Counter
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
TAREFAS_SRC = ROOT / "docs" / "tarefas.md"
INCIDENTES_DIR = ROOT / "docs" / "incidentes"
PORTAL_RELATORIOS = ROOT / "portal" / "docs" / "relatorios"
PORTAL_INDEX = ROOT / "portal" / "docs" / "index.md"
REPORTS_DIR = ROOT / "reports"
MKDOCS_YML = ROOT / "portal" / "mkdocs.yml"

MESES_PT = {
    1: "Janeiro", 2: "Fevereiro", 3: "Março", 4: "Abril",
    5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
    9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro",
}


def parse_tasks_for_month(text, year, month):
    concluidas = []
    abertas_andamento = []
    in_concluidas = False

    for line in text.splitlines():
        stripped = line.strip()

        if stripped.startswith("## Abertas") or stripped.startswith("## Em andamento"):
            in_concluidas = False
        elif stripped.startswith("## Concluídas"):
            in_concluidas = True

        if not stripped.startswith("- ["):
            continue

        task = parse_task(stripped)
        if not task:
            continue

        if in_concluidas and task["date"]:
            dt = datetime.strptime(task["date"], "%Y-%m-%d")
            if dt.year == year and dt.month == month:
                concluidas.append(task)
        elif not in_concluidas:
            if task["status"] in ("open", "in_progress"):
                abertas_andamento.append(task)

    return concluidas, abertas_andamento


def parse_task(line):
    checkbox_match = re.match(r"- \[(.)\] ", line)
    if not checkbox_match:
        return None

    mark = checkbox_match.group(1)
    status = {"x": "done", "~": "in_progress", " ": "open"}.get(mark, "open")
    rest = line[len(checkbox_match.group(0)):]

    task_id = None
    id_match = re.match(r"(T\d+)\s+", rest)
    if id_match:
        task_id = id_match.group(1)
        rest = rest[len(id_match.group(0)):]

    cat_match = re.match(r"\[([^\]]+)\]\s*", rest)
    category = "Shift LIS"
    if cat_match:
        category = cat_match.group(1)
        rest = rest[len(cat_match.group(0)):]

    date = None
    date_match = re.search(r"\*(?:concluída|iniciada|aberta) em (\d{4}-\d{2}-\d{2})\*", rest)
    if date_match:
        date = date_match.group(1)

    desc = rest
    desc = re.sub(r"\s*—\s*\*(?:concluída|iniciada|aberta) em \d{4}-\d{2}-\d{2}\*", "", desc)
    desc = re.sub(r"\s*\(agendado para [^)]+\)", "", desc)
    desc = re.sub(r"\s*\(resp\. [^)]+\)", "", desc)
    desc = desc.strip().rstrip(" —").strip()

    return {
        "status": status,
        "category": category,
        "description": desc,
        "date": date,
    }


def get_week_number(date_str, month):
    dt = datetime.strptime(date_str, "%Y-%m-%d")
    first = datetime(dt.year, month, 1)
    return ((dt.day - 1) // 7) + 1


def get_week_range(year, month, week_num):
    start_day = (week_num - 1) * 7 + 1
    end_day = min(start_day + 6, 28 if month == 2 else 30 if month in (4, 6, 9, 11) else 31)
    try:
        datetime(year, month, end_day)
    except ValueError:
        end_day = 28 if month == 2 else 30
    return f"{start_day:02d}/{month:02d}", f"{end_day:02d}/{month:02d}"


def find_incidents(year, month):
    incidents = []
    if not INCIDENTES_DIR.exists():
        return incidents
    prefix = f"{year}-{month:02d}"
    for f in sorted(INCIDENTES_DIR.iterdir()):
        if f.name.startswith(prefix) and f.suffix == ".md":
            content = f.read_text(encoding="utf-8")
            title_match = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
            title = title_match.group(1) if title_match else f.stem
            incidents.append({"file": f.name, "title": title, "date": f.name[:10]})
    return incidents


def generate_md_report(year, month, concluidas, pendentes, incidents):
    month_name = MESES_PT[month]
    today = datetime.now()
    is_current = (year == today.year and month == today.month)
    period_end = f"{today.day:02d}/{month:02d}" if is_current else f"{28 if month == 2 else 30 if month in (4,6,9,11) else 31}/{month:02d}"
    status_note = " (parcial — mês em andamento)" if is_current else ""

    lines = [
        f"# Resumo Mensal — {month_name} {year}",
        "",
        f"**Período:** 01/{month:02d} a {period_end}{status_note}",
        "",
        "---",
        "",
        f"## Tarefas concluídas: {len(concluidas)}",
        "",
        "| # | Categoria | Descrição | Data |",
        "|---|-----------|-----------|------|",
    ]

    for i, t in enumerate(sorted(concluidas, key=lambda x: x["date"]), 1):
        dt = datetime.strptime(t["date"], "%Y-%m-%d")
        lines.append(f"| {i} | {t['category']} | {t['description']} | {dt.day:02d}/{dt.month:02d} |")

    cat_counts = Counter(t["category"] for t in concluidas)
    lines.extend([
        "",
        "## Por categoria",
        "",
        "| Categoria | Qtd |",
        "|-----------|-----|",
    ])
    for cat, count in cat_counts.most_common():
        lines.append(f"| {cat} | {count} |")
    lines.append(f"| **Total** | **{len(concluidas)}** |")

    weeks = defaultdict(int)
    for t in concluidas:
        wn = get_week_number(t["date"], month)
        weeks[wn] += 1

    lines.extend([
        "",
        "## Distribuição semanal",
        "",
        "| Semana | Concluídas |",
        "|--------|-----------|",
    ])
    for wn in sorted(weeks.keys()):
        start, end = get_week_range(year, month, wn)
        lines.append(f"| Semana {wn} ({start}–{end}) | {weeks[wn]} |")

    if pendentes:
        lines.extend(["", "## Tarefas em andamento / pendentes", ""])
        for t in pendentes:
            lines.append(f"- {t['description']} — {t['category']}")

    lines.extend(["", "## Incidentes registrados", ""])
    if incidents:
        for inc in incidents:
            lines.append(f"- **{inc['date']}** — {inc['title']}")
    else:
        lines.append(f"Nenhum incidente registrado em {month_name.lower()}.")

    lines.append("")
    return "\n".join(lines)


CATEGORY_COLORS = {
    "Shift LIS": "00897B",
    "Shift Integração": "00ACC1",
    "Shift Automação": "5C6BC0",
    "Rede/Servidores": "546E7A",
    "Municípios": "FFA000",
    "Shift Cadastro": "8D6E63",
}
DEFAULT_CATEGORY_COLOR = "78909C"
AMBER = "F9A825"
RED = "C62828"
DARK_HEX = "212121"
GRAY_HEX = "757575"


def cat_color(category):
    return CATEGORY_COLORS.get(category, DEFAULT_CATEGORY_COLOR)


def fit_size(n_items):
    """Escolhe tamanho de fonte para uma lista com marcadores, evitando overflow."""
    for threshold, size in ((6, 18), (10, 16), (16, 14), (24, 12)):
        if n_items <= threshold:
            return size
    return 11


def replace_text_in_slide(slide, replacements):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_fill_color(slide, shape_name, color_hex):
    from pptx.dml.color import RGBColor
    shape = find_shape(slide, shape_name)
    if shape is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)


def _set_bullet_char(paragraph, char, color_hex):
    from pptx.oxml.ns import qn
    pPr = paragraph._p.get_or_add_pPr()
    buClr = pPr.makeelement(qn("a:buClr"), {})
    buClr.append(buClr.makeelement(qn("a:srgbClr"), {"val": color_hex}))
    buSzPct = pPr.makeelement(qn("a:buSzPct"), {"val": "65000"})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    defRPr = pPr.find(qn("a:defRPr"))
    for el in (buClr, buSzPct, buFont, buChar):
        if defRPr is not None:
            defRPr.addprevious(el)
        else:
            pPr.append(el)


def add_bullet_paragraph(text_frame, idx, bullet_color, space_after_pt):
    from pptx.util import Pt
    paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
    paragraph.space_after = Pt(space_after_pt)
    paragraph.line_spacing = 1.15
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "228600")
    pPr.set("indent", "-228600")
    _set_bullet_char(paragraph, "●", bullet_color)
    return paragraph


def add_run(paragraph, text, size_pt, color_hex, bold=False, italic=False):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor.from_string(color_hex)
    return run


def write_message(text_frame, message):
    """Escreve uma única linha sem marcador, para estados vazios (ex.: sem incidentes)."""
    paragraph = text_frame.paragraphs[0]
    add_run(paragraph, message, 16, GRAY_HEX, italic=True)


def prep_text_frame(text_frame):
    """Impede que a caixa de texto cresça além dos limites do slide (o padrão do
    python-pptx é SHAPE_TO_FIT_TEXT, que 'empurra' o conteúdo pra fora do slide)."""
    from pptx.enum.text import MSO_AUTO_SIZE
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE


def populate_task_bullets(text_frame, tasks, bullet_color, font_size, space_after):
    prep_text_frame(text_frame)
    for i, t in enumerate(tasks):
        dt = datetime.strptime(t["date"], "%Y-%m-%d")
        paragraph = add_bullet_paragraph(text_frame, i, bullet_color, space_after)
        add_run(paragraph, t["description"], font_size, DARK_HEX)
        add_run(paragraph, f"   {dt.day:02d}/{dt.month:02d}", font_size - 2, GRAY_HEX, italic=True)


def clone_slide(prs, template_slide):
    """Clone a slide from the template into the presentation."""
    import io
    from copy import deepcopy
    from pptx.oxml.ns import qn

    slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in template_slide.shapes:
        el = deepcopy(shape._element)
        if shape.shape_type == 13:  # PICTURE — precisa recriar a relação da imagem
            blip = el.find(f".//{qn('a:blip')}")
            if blip is not None:
                _, new_rid = new_slide.part.get_or_add_image_part(io.BytesIO(shape.image.blob))
                blip.set(qn("r:embed"), new_rid)
        new_slide.shapes._spTree.append(el)

    return new_slide


def generate_pptx(year, month, concluidas, pendentes, incidents):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    TEMPLATE_PATH = ROOT / "reports" / "template.pptx"

    if not TEMPLATE_PATH.exists():
        print(f"  Aviso: template não encontrado em {TEMPLATE_PATH}, gerando sem template")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs

    template = Presentation(str(TEMPLATE_PATH))
    prs = Presentation()
    prs.slide_width = template.slide_width
    prs.slide_height = template.slide_height

    month_name = MESES_PT[month]
    titulo_mes = f"{month_name} {year}"
    cat_counts = Counter(t["category"] for t in concluidas)
    cats_sorted = cat_counts.most_common()

    # --- Slide 0: Capa ---
    slide = clone_slide(prs, template.slides[0])
    replace_text_in_slide(slide, {
        "{TITULO_MES}": titulo_mes,
        "{TOTAL_CONCLUIDAS}": str(len(concluidas)),
        "{TOTAL_PENDENTES}": str(len(pendentes)),
        "{TOTAL_INCIDENTES}": str(len(incidents)),
    })
    for shape_name, color_hex in (
        ("TextBox 10", "00897B"), ("TextBox 13", AMBER), ("TextBox 16", RED),
    ):
        shape = find_shape(slide, shape_name)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor.from_string(color_hex)

    # --- Slide 1: Resumo por categoria (gráfico de barras) ---
    slide = clone_slide(prs, template.slides[1])
    replace_text_in_slide(slide, {"{TITULO_MES}": titulo_mes})
    placeholder = find_shape(slide, "TextBox 5")
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    placeholder._element.getparent().remove(placeholder._element)

    subtitle = slide.shapes.add_textbox(left, top - Pt(28), width, Pt(28))
    add_run(subtitle.text_frame.paragraphs[0], f"Total: {len(concluidas)} tarefas concluídas", 14, GRAY_HEX)

    chart_data = CategoryChartData()
    chart_data.categories = [c for c, _ in cats_sorted]
    chart_data.add_series("Tarefas concluídas", [n for _, n in cats_sorted])
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.category_axis.reverse_order = True
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor.from_string(DARK_HEX)
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(14)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = RGBColor.from_string(DARK_HEX)
    series = chart.series[0]
    for i, (cat, _) in enumerate(cats_sorted):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(cat_color(cat))

    # --- Slides por categoria (clona slide 3 do template) ---
    COLUMN_SPLIT_THRESHOLD = 10
    COLUMN_GAP = 365760  # ~0.4"

    for cat, count in cats_sorted:
        slide = clone_slide(prs, template.slides[3])
        set_fill_color(slide, "Rectangle 1", cat_color(cat))
        replace_text_in_slide(slide, {
            "{CATEGORIA_NOME}": cat,
            "{CATEGORIA_QTD}": str(count),
            "{TITULO_MES}": titulo_mes,
        })
        cat_tasks = sorted(
            [t for t in concluidas if t["category"] == cat],
            key=lambda x: x["date"]
        )
        placeholder = find_shape(slide, "TextBox 5")
        left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

        if len(cat_tasks) > COLUMN_SPLIT_THRESHOLD:
            placeholder._element.getparent().remove(placeholder._element)
            col_width = (width - COLUMN_GAP) // 2
            mid = (len(cat_tasks) + 1) // 2
            columns = [cat_tasks[:mid], cat_tasks[mid:]]
            size = fit_size(mid)
            space_after = max(4, size - 6)
            for col_idx, col_tasks in enumerate(columns):
                col_left = left + col_idx * (col_width + COLUMN_GAP)
                box = slide.shapes.add_textbox(col_left, top, col_width, height)
                populate_task_bullets(box.text_frame, col_tasks, cat_color(cat), size, space_after)
        else:
            text_frame = placeholder.text_frame
            text_frame.clear()
            size = fit_size(len(cat_tasks))
            space_after = max(4, size - 6)
            populate_task_bullets(text_frame, cat_tasks, cat_color(cat), size, space_after)

    # --- Slide: Pendências (clona slide 4) ---
    if pendentes:
        slide = clone_slide(prs, template.slides[4])
        set_fill_color(slide, "Rectangle 1", AMBER)
        replace_text_in_slide(slide, {"{TITULO_MES}": titulo_mes})
        text_frame = find_shape(slide, "TextBox 5").text_frame
        text_frame.clear()
        prep_text_frame(text_frame)
        size = fit_size(len(pendentes))
        space_after = max(4, size - 6)
        for i, t in enumerate(pendentes):
            paragraph = add_bullet_paragraph(text_frame, i, cat_color(t["category"]), space_after)
            add_run(paragraph, f"[{t['category']}] ", size, cat_color(t["category"]), bold=True)
            add_run(paragraph, t["description"], size, DARK_HEX)

    # --- Slide: Incidentes (clona slide 5) ---
    slide = clone_slide(prs, template.slides[5])
    set_fill_color(slide, "Rectangle 1", RED)
    replace_text_in_slide(slide, {"{TITULO_MES}": titulo_mes})
    text_frame = find_shape(slide, "TextBox 5").text_frame
    text_frame.clear()
    prep_text_frame(text_frame)
    if incidents:
        size = fit_size(len(incidents))
        space_after = max(4, size - 6)
        for i, inc in enumerate(incidents):
            paragraph = add_bullet_paragraph(text_frame, i, RED, space_after)
            add_run(paragraph, f"{inc['date']}  ", size, RED, bold=True)
            add_run(paragraph, inc["title"], size, DARK_HEX)
    else:
        write_message(text_frame, f"Nenhum incidente registrado em {month_name.lower()}.")

    return prs


def update_relatorios_index(year, month):
    index_path = PORTAL_RELATORIOS / "index.md"
    if not index_path.exists():
        return

    content = index_path.read_text(encoding="utf-8")
    month_name = MESES_PT[month]
    entry = f"{year}-{month:02d}"

    if entry in content:
        return

    today = datetime.now()
    is_current = (year == today.year and month == today.month)
    status = ":material-clock-outline: Parcial (mês em andamento)" if is_current else ":material-check-circle: Completo"

    new_row = f"| [{month_name} {year}]({entry}.md) | {status} |"

    table_end = content.rfind("|")
    if table_end != -1:
        line_end = content.find("\n", table_end)
        if line_end == -1:
            line_end = len(content)
        content = content[:line_end] + "\n" + new_row + content[line_end:]

    index_path.write_text(content, encoding="utf-8")
    print(f"  Atualizado: relatorios/index.md")


def update_portal_index(concluidas_count, month_name):
    if not PORTAL_INDEX.exists():
        return

    content = PORTAL_INDEX.read_text(encoding="utf-8")

    content = re.sub(
        r'(Concluídas em \w+\s*\|\s*)\d+',
        f"Concluídas em {month_name.lower()} | {concluidas_count}",
        content
    )

    today = datetime.now()
    content = re.sub(
        r'Atualizado em \d{2}/\d{2}/\d{4}',
        f"Atualizado em {today.day:02d}/{today.month:02d}/{today.year}",
        content
    )

    PORTAL_INDEX.write_text(content, encoding="utf-8")
    print(f"  Atualizado: portal index.md")


def update_mkdocs_nav(year, month):
    content = MKDOCS_YML.read_text(encoding="utf-8")
    month_name = MESES_PT[month]
    entry = f'"{month_name} {year}": relatorios/{year}-{month:02d}.md'

    if entry in content:
        return

    relatorios_start = content.find("  - Relatórios:")
    if relatorios_start == -1:
        return

    next_section = content.find("\n  - ", relatorios_start + 1)
    if next_section == -1:
        next_section = len(content)

    insert_pos = content.find("\n", content.find("relatorios/index.md", relatorios_start))
    if insert_pos == -1 or insert_pos > next_section:
        return

    new_line = f'\n    - "{month_name} {year}": relatorios/{year}-{month:02d}.md'
    content = content[:insert_pos] + new_line + content[insert_pos:]

    MKDOCS_YML.write_text(content, encoding="utf-8")
    print(f"  Atualizado: mkdocs.yml nav (relatórios)")


def main():
    parser = argparse.ArgumentParser(description="Gera relatório mensal")
    parser.add_argument("--mes", type=str, default=None, help="Mês no formato YYYY-MM (padrão: mês atual)")
    args = parser.parse_args()

    if args.mes:
        year, month = map(int, args.mes.split("-"))
    else:
        today = datetime.now()
        year, month = today.year, today.month

    month_name = MESES_PT[month]
    print(f"Gerando relatório: {month_name} {year}")

    if not TAREFAS_SRC.exists():
        print(f"Erro: {TAREFAS_SRC} não encontrado")
        return

    text = TAREFAS_SRC.read_text(encoding="utf-8")
    concluidas, pendentes = parse_tasks_for_month(text, year, month)
    incidents = find_incidents(year, month)

    print(f"  {len(concluidas)} tarefas concluídas, {len(pendentes)} pendentes, {len(incidents)} incidentes")

    # Gera MD para portal
    PORTAL_RELATORIOS.mkdir(parents=True, exist_ok=True)
    md_content = generate_md_report(year, month, concluidas, pendentes, incidents)
    md_path = PORTAL_RELATORIOS / f"{year}-{month:02d}.md"
    md_path.write_text(md_content, encoding="utf-8")
    print(f"  Gerado: {md_path.relative_to(ROOT)}")

    # Gera PPTX para reports/
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    pptx = generate_pptx(year, month, concluidas, pendentes, incidents)
    pptx_path = REPORTS_DIR / f"{year}-{month:02d}_resumo-mensal.pptx"
    pptx.save(str(pptx_path))
    print(f"  Gerado: {pptx_path.relative_to(ROOT)}")

    # Atualiza índices
    update_relatorios_index(year, month)
    update_portal_index(len(concluidas), month_name)
    update_mkdocs_nav(year, month)

    print(f"\nRelatório de {month_name} {year} concluído.")


if __name__ == "__main__":
    main()
