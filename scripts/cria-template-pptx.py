"""
Cria o template PPTX com branding FUNFARME para relatórios mensais.
Executar uma vez — o template fica em reports/template.pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "reports" / "template.pptx"
ASSETS = ROOT / "reports" / "assets"

TEAL = RGBColor(0x00, 0x89, 0x7B)
TEAL_DARK = RGBColor(0x00, 0x5F, 0x56)
TEAL_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)
CYAN = RGBColor(0x00, 0xBC, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

W = Inches(13.333)
H = Inches(7.5)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return box


def add_logo(slide, path, left, top, height):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, height=height)


def add_header(slide, title, dark=False):
    bg_color = TEAL_DARK if dark else TEAL
    add_shape(slide, 0, 0, W, Inches(1.1), bg_color)
    add_shape(slide, 0, Inches(1.1), W, Inches(0.04), CYAN)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
             title, size=28, color=WHITE, bold=True)
    # NTI logo no header (direita)
    nti_path = ASSETS / "nti.png"
    if nti_path.exists():
        slide.shapes.add_picture(str(nti_path), Inches(11.2), Inches(0.2), height=Inches(0.65))


def add_footer(slide):
    add_shape(slide, 0, Inches(6.9), W, Inches(0.6), LIGHT_GRAY)
    add_text(slide, Inches(0.8), Inches(6.95), Inches(5), Inches(0.4),
             "{TITULO_MES}", size=10, color=GRAY)
    add_text(slide, Inches(7), Inches(6.95), Inches(5.5), Inches(0.4),
             "TI Laboratório — FUNFARME", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def build_template():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    funfarme_logo = ASSETS / "funfarme.jpg"
    hb_logo = ASSETS / "hospital-de-base.png"

    # ===================== SLIDE 1: CAPA =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Faixa teal (topo)
    add_shape(slide, 0, 0, W, Inches(3.5), TEAL)
    # Linha accent
    add_shape(slide, 0, Inches(3.5), W, Inches(0.06), CYAN)

    # Logos no canto superior direito da faixa teal
    add_logo(slide, funfarme_logo, Inches(10.2), Inches(0.3), Inches(1.1))
    add_logo(slide, hb_logo, Inches(11.6), Inches(0.35), Inches(0.9))

    # Título
    add_text(slide, Inches(1), Inches(1.0), Inches(8.5), Inches(1),
             "RELATÓRIO MENSAL", size=44, color=WHITE, bold=True)

    # Subtítulo
    add_text(slide, Inches(1), Inches(2.0), Inches(8.5), Inches(0.6),
             "{TITULO_MES}", size=26, color=TEAL_LIGHT)

    # Instituição
    add_text(slide, Inches(1), Inches(2.6), Inches(8.5), Inches(0.5),
             "TI Laboratório Central — FUNFARME", size=16, color=TEAL_LIGHT)

    # Métricas (parte branca)
    # Card 1: Concluídas
    add_shape(slide, Inches(1), Inches(4.2), Inches(3.2), Inches(2.0), LIGHT_GRAY)
    add_text(slide, Inches(1), Inches(4.35), Inches(3.2), Inches(0.4),
             "Tarefas concluídas", size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4.8), Inches(3.2), Inches(1.0),
             "{TOTAL_CONCLUIDAS}", size=52, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Card 2: Pendentes
    add_shape(slide, Inches(5.1), Inches(4.2), Inches(3.2), Inches(2.0), LIGHT_GRAY)
    add_text(slide, Inches(5.1), Inches(4.35), Inches(3.2), Inches(0.4),
             "Pendentes", size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.1), Inches(4.8), Inches(3.2), Inches(1.0),
             "{TOTAL_PENDENTES}", size=52, color=DARK, bold=True, align=PP_ALIGN.CENTER)

    # Card 3: Incidentes
    add_shape(slide, Inches(9.2), Inches(4.2), Inches(3.2), Inches(2.0), LIGHT_GRAY)
    add_text(slide, Inches(9.2), Inches(4.35), Inches(3.2), Inches(0.4),
             "Incidentes", size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(9.2), Inches(4.8), Inches(3.2), Inches(1.0),
             "{TOTAL_INCIDENTES}", size=52, color=DARK, bold=True, align=PP_ALIGN.CENTER)

    # Rodapé da capa
    add_shape(slide, 0, Inches(6.9), W, Inches(0.6), LIGHT_GRAY)
    add_text(slide, Inches(1), Inches(6.95), Inches(11), Inches(0.4),
             "FUNFARME — Fundação Faculdade Regional de Medicina · São José do Rio Preto — SP",
             size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # ===================== SLIDE 2: RESUMO POR CATEGORIA =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Resumo por Categoria")
    add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
             "{CATEGORIAS}", size=16, color=DARK)
    add_footer(slide)

    # ===================== SLIDE 3: DISTRIBUIÇÃO SEMANAL =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Distribuição Semanal")
    add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
             "{SEMANAS}", size=16, color=DARK)
    add_footer(slide)

    # ===================== SLIDE 4: DETALHE CATEGORIA =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "{CATEGORIA_NOME} ({CATEGORIA_QTD})")
    add_text(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5),
             "{CATEGORIA_ITENS}", size=14, color=DARK)
    add_footer(slide)

    # ===================== SLIDE 5: PENDÊNCIAS =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Pendências", dark=True)
    add_text(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5),
             "{PENDENCIAS}", size=14, color=DARK)
    add_footer(slide)

    # ===================== SLIDE 6: INCIDENTES =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Incidentes", dark=True)
    add_text(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5),
             "{INCIDENTES}", size=14, color=DARK)
    add_footer(slide)

    prs.save(str(OUT))
    print(f"Template salvo: {OUT}")


if __name__ == "__main__":
    build_template()
